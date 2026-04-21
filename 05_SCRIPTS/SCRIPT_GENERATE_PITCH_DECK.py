import os
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


def create_pitch_deck():
    prs = Presentation()
    repo_root = Path(__file__).resolve().parents[1]
    # Coloque PNGs opcionais em 06_DATA/pitch_deck_assets/ com estes nomes; senão o slide segue só com texto.
    IMG_DIR = repo_root / "06_DATA" / "pitch_deck_assets"
    IMG_COVER = IMG_DIR / "fluxus_pitch_cover_apple_style_1776265575631.png"
    IMG_NETWORK = IMG_DIR / "fluxus_global_network_apple_style_1776265600098.png"
    IMG_UI = IMG_DIR / "fluxus_cockpit_ui_apple_style_1776265623017.png"
    IMG_ATOMIC = IMG_DIR / "fluxus_atomic_settlement_apple_style_1776265642427.png"

    def add_image_slide(prs, img_path, title_pt, title_en, body_pt, body_en):
        slide = prs.slides.add_slide(prs.slide_layouts[6]) # Blank layout
        
        # Add background image if exists
        pimg = str(img_path) if img_path else ""
        if pimg and os.path.exists(pimg):
            prs.slide_width = Inches(13.33)
            prs.slide_height = Inches(7.5)
            slide.shapes.add_picture(pimg, 0, 0, width=prs.slide_width, height=prs.slide_height)

        # Add Overlay for text readability
        # (This is a bit complex in pptx without manual shapes, let's just add text boxes)
        
        # Title Box
        txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12), Inches(1.5))
        tf = txBox.text_frame
        p = tf.add_paragraph()
        p.text = f"{title_pt}\n{title_en}"
        p.font.bold = True
        p.font.size = Pt(36)
        p.font.name = 'Arial'
        p.font.color.rgb = RGBColor(255, 255, 255) # White

        # Body Box
        txBox2 = slide.shapes.add_textbox(Inches(0.5), Inches(5.5), Inches(12), Inches(1.5))
        tf2 = txBox2.text_frame
        p2 = tf2.add_paragraph()
        p2.text = f"{body_pt}\n{body_en}"
        p2.font.size = Pt(18)
        p2.font.name = 'Arial'
        p2.font.color.rgb = RGBColor(200, 200, 200) # Light Grey

    # Slide 1: Cover
    add_image_slide(prs, IMG_COVER, 
                    "FLUXUS: O Nexus da Liquidez Global", "FLUXUS: The Nexus of Global Liquidity",
                    "O Protocolo de Rede da Transparência Financeira", "The Network Protocol of Financial Transparency")

    # Slide 2: The Problem
    add_image_slide(prs, IMG_NETWORK, 
                    "A Fricção Invisível do Mundo", "The World's Invisible Friction",
                    "$669 Bilhões evaporam anualmente em taxas bancárias obsoletas.", "$669 Billion evaporate annually in obsolete banking fees.")

    # Slide 3: Silo vs Protocol
    add_image_slide(prs, "", 
                    "Dos Silos ao TCP/IP do Dinheiro", "From Silos to the TCP/IP of Money",
                    "O FLUXUS orquestra a liquidez horizontalmente.", "FLUXUS orchestrates liquidity horizontally.")

    # Slide 4: Atomic Settlement
    add_image_slide(prs, IMG_ATOMIC, 
                    "Liquidação Atômica (T+0)", "Atomic Settlement (T+0)",
                    "Settlement síncrono em < 15 segundos.", "Synchronous settlement in < 15 seconds.")

    # Slide 5: UI Mockup
    add_image_slide(prs, IMG_UI, 
                    "A Matemática do Unicórnio", "The Unicorn Mathematics",
                    "GTV: $17.5B | Lucro Líquido: $189.9M (60 Meses).", "GTV: $17.5B | Net Profit: $189.9M (60 Months).")

    # Final Save
    OUTPUT_FILE = repo_root / "02_INVESTOR" / "INVESTOR_PITCH_DECK_AUTOGEN.pptx"
    prs.save(str(OUTPUT_FILE))
    print(f"Pitch Deck salvo em: {OUTPUT_FILE.resolve()}")

if __name__ == "__main__":
    create_pitch_deck()
